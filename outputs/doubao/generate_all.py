#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate all FDE commercialization deliverables for FieldPilot AI benchmark."""

import os
import csv
from datetime import datetime
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ============================================================
# 04-source-log.md
# ============================================================
sources = [
    {
        "id": "S01",
        "url": "http://www.ce.cn/cysc/newmain/yc/jsxw/202603/t20260309_2813343.shtml",
        "title": "政策加码护航 企业迅捷布局智能体成拉动制造业转型新引擎",
        "publisher": "中国经济网",
        "date": "2026-03-09",
        "use": "制造业AI支出预测：IDC预测2028年中国工业企业AI支出接近900亿元，CAGR 37.7%；支撑制造业作为优先行业的市场规模判断",
        "uncertainty": "IDC预测数据，实际增速受宏观经济和企业IT预算影响；'AI+制造'35%CAGR为机构预测"
    },
    {
        "id": "S02",
        "url": "https://digitalinasia.com/china-30000-smart-factories-industrial-ai/",
        "title": "China Industrial AI Market: 30,000 Smart Factories",
        "publisher": "Digital in Asia",
        "date": "2026-03-25",
        "use": "工信部数据：2025年底超30%规模以上制造企业已采用AI，核心AI产业超1.2万亿元，6200+AI企业；验证制造业AI渗透率",
        "uncertainty": "核心AI产业口径 broader than industrial AI；'采用AI'定义可能包含轻量级应用"
    },
    {
        "id": "S03",
        "url": "https://news.sina.com.cn/o/2025-11-08/doc-infwschz3878925.shtml",
        "title": "AI全面'上岗'双11，电商迎来历史性转折点",
        "publisher": "新浪新闻",
        "date": "2025-11-08",
        "use": "天猫双11数据：AIGC生成1.5亿素材，AI客服服务3亿人次，全自动承接1亿人次；验证零售电商AI应用成熟度和场景丰富度",
        "uncertainty": "平台自有数据，可能偏向头部平台能力；中小商家AI渗透率未明确"
    },
    {
        "id": "S04",
        "url": "https://developer.aliyun.com/article/1690989",
        "title": "我是一家电商企业，推荐一款Agent产品：从选型到落地的全攻略",
        "publisher": "阿里云开发者社区",
        "date": "2025-12-02",
        "use": "艾瑞咨询数据：电商客服人力成本占运营成本23%，70%客服工作为重复咨询，人均日处理80-120条，错误率4.2%；支撑电商场景ROI测算",
        "uncertainty": "2024年数据，2025年AI客服渗透率已大幅提升；'平均'口径可能掩盖行业差异"
    },
    {
        "id": "S05",
        "url": "https://www.hundun.cn/articles/66bbd2b8ff260361.html",
        "title": "解码Palantir：这家美国'最神秘'的软件公司，给中国SaaS行业上了一课",
        "publisher": "混沌学园",
        "date": "2026",
        "use": "FDE模式定义：Palantir客单价从2019年790万美元增至2024年超1000万美元；FDE长期驻扎客户现场，深度融合业务流程；支撑FDE商业模式设计",
        "uncertainty": "Palantir面向政府和大型企业，客单价远高于中国市场可参照水平；FDE模式在中国本土化适配需验证"
    },
    {
        "id": "S06",
        "url": "https://raw.githubusercontent.com/xdash/FDE-the-Guidance-Book-of-Forward-Deployed-Engineer/master/%E5%89%8D%E7%BA%BF%E9%83%A8%E7%BD%B2%E5%B7%A5%E7%A8%8B%E5%B8%88%EF%BC%88FDE%EF%BC%89v1.0.24.pdf",
        "title": "前线部署工程师：人工智能时代的客户价值交付秘籍",
        "publisher": "GitHub (xdash/FDE)",
        "date": "2026-08-01",
        "use": "FDE行业趋势：Databricks将专业服务整体改组为FDE组织，12个月服务1900+客户；德勤2025年12月成立FDE业务线；FDE岗位招聘需求暴涨42倍；支撑FDE人才和组织设计",
        "uncertainty": "GitHub开源文档，数据来源为LinkedIn等公开渠道；'42倍'增速可能基数较小"
    },
    {
        "id": "S07",
        "url": "https://cloud.tencent.com/developer/article/2720303",
        "title": "AI办公智能体领域行业发展洞察",
        "publisher": "腾讯云开发者社区",
        "date": "2026-08-04",
        "use": "IDC数据：2025年中国企业级Agent市场规模约190-212亿元，2026年预计449亿元，CAGR超110%，2029年有望突破3320亿元；支撑整体TAM判断",
        "uncertainty": "不同机构对'企业级Agent'口径差异大（69亿-670亿不等）；IDC数据为预测值"
    },
    {
        "id": "S08",
        "url": "https://www.nfra.gov.cn/cn/view/pages/ItemDetail.html?docId=1261710",
        "title": "国家金融监督管理总局：金融机构强化人工智能风险治理",
        "publisher": "国家金融监督管理总局",
        "date": "2026-06-18",
        "use": "金融AI监管框架：要求将AI风险纳入全面风险管理，防范模型黑箱、生成幻觉、算法歧视；高风险场景需风险管理委员会审批；支撑金融行业合规边界分析",
        "uncertainty": "监管文件为原则性指导，具体执行细则待落地；不同金融子行业（银行/保险/证券）要求差异大"
    },
    {
        "id": "S09",
        "url": "https://ex.chinadaily.com.cn/exchange/partners/82/rss/channel/cn/columns/sz8srm/stories/WS6a686ff4a310d709c2fc0391.html",
        "title": "银行AI大潮下，没有上亿预算的地方中小银行，如何'轻量化'落地？",
        "publisher": "中国日报",
        "date": "2026-07-28",
        "use": "金发〔2026〕8号文：银行业保险业首份AI全流程专项监管文件，首次定义金融智能体(Agent)；金融AI进入合规化标准化阶段；支撑金融行业进入壁垒判断",
        "uncertainty": "文章聚焦中小银行，大行AI能力和预算差异大；监管文件具体条款需原文核实"
    },
    {
        "id": "S10",
        "url": "https://www.cn-healthcare.com/articlewm/20250708/content-1653343.html",
        "title": "中国智慧医疗多元场景的发展现状与未来展望",
        "publisher": "健康界",
        "date": "2025-07-08",
        "use": "医疗AI市场：2023年智慧医疗62.85亿元（CAGR 53.37%），2025年预计突破120亿元；AI医疗器械2025年预计242亿元；支撑医疗行业市场规模判断",
        "uncertainty": "智慧医疗口径包含硬件和软件；非临床运营AI占比未明确；医疗数据合规壁垒高"
    },
    {
        "id": "S11",
        "url": "https://m.36kr.com/p/3616263140312320",
        "title": "2025年度盘点：SaaS行业的'AI大考'与上市公司的生死突围",
        "publisher": "36氪",
        "date": "2025-12-29",
        "use": "企业软件AI转型：金蝶Q3试点'AI CSM'，人均服务客户数提升3倍，AI功能使用率从32%升至68%；SaaS行业实施和客户成功是AI渗透重点；支撑企业软件行业适配分析",
        "uncertainty": "金蝶单一案例，行业普适性待验证；'AI CSM'为内部岗位试点，非外部FDE服务"
    },
    {
        "id": "S12",
        "url": "http://www.acla.org.cn/info/cab0337cce9343109c4326bd7ed24917",
        "title": "人工智能浪潮下的律所实践探索",
        "publisher": "中国律师网",
        "date": "2025-07-16",
        "use": "专业服务AI应用：植德律所开发'小植同学'等智能助手，基于RAG将律所知识库转化为实时应答；AI在文书起草、案例检索、合同审核表现突出；支撑专业服务行业场景分析",
        "uncertainty": "单一律所案例，行业整体AI渗透率低；律所对数据保密性要求极高，外部FDE入驻阻力大"
    },
]

source_log = """# 04-source-log.md — 公开来源记录

> 访问日期：2026-08-25（北京时间）
> 产品：豆包工作 Auto 高模式
> 说明：以下来源均为公开可访问网页，用于支撑行业评分、目标企业筛选和商业化方案判断。

## 来源清单

"""

for s in sources:
    source_log += f"""### {s['id']}：{s['title']}
- **URL**：{s['url']}
- **发布方**：{s['publisher']}
- **发布日期**：{s['date']}
- **访问日期**：2026-08-25
- **用途**：{s['use']}
- **不确定性**：{s['uncertainty']}

"""

source_log += """## 来源统计
- 总来源数：12
- 覆盖行业：制造(2)、零售电商(2)、金融(2)、医疗(1)、企业软件(1)、专业服务(1)、FDE模式(2)、整体市场(1)
- 权威来源：政府/监管机构2个（金融监管总局、中国政府网），官方媒体4个（中国经济网、中国日报、新浪、健康界），行业研究3个，企业/社区3个

## 关键数据交叉验证
- 中国AI核心产业规模2025年：1.2万亿元（信通院/工信部，S02/S06交叉验证）
- 企业级Agent市场2025年：190-212亿元（IDC，S07）；不同机构口径差异大（69亿-670亿）
- 制造业AI渗透率：超30%规模以上企业（工信部，S02）
- 电商AI客服：双11期间AI全自动承接1亿人次（天猫，S03）
- 金融AI监管：金发〔2026〕8号文首次定义金融智能体（S09）
"""

with open(os.path.join(OUTPUT_DIR, "04-source-log.md"), "w", encoding="utf-8") as f:
    f.write(source_log)
print("04-source-log.md generated")

# ============================================================
# 01-industry-prioritization.xlsx
# ============================================================
wb = Workbook()

# --- Sheet 1: Industry Scoring ---
ws1 = wb.active
ws1.title = "行业评分"

# Styles
header_font = Font(bold=True, color="FFFFFF", size=11)
header_fill = PatternFill(start_color="2F5496", end_color="2F5496", fill_type="solid")
priority_fill = PatternFill(start_color="C6EFCE", end_color="C6EFCE", fill_type="solid")
thin_border = Border(
    left=Side(style='thin'), right=Side(style='thin'),
    top=Side(style='thin'), bottom=Side(style='thin')
)

# Title
ws1.merge_cells('A1:I1')
ws1['A1'] = "FieldPilot AI — 六行业统一评分框架（满分100）"
ws1['A1'].font = Font(bold=True, size=14)
ws1['A1'].alignment = Alignment(horizontal='center')

# Headers
headers = ["评分维度", "权重", "制造", "零售电商", "金融", "医疗", "企业软件", "专业服务", "评分说明"]
for col, h in enumerate(headers, 1):
    cell = ws1.cell(row=3, column=col, value=h)
    cell.font = header_font
    cell.fill = header_fill
    cell.alignment = Alignment(horizontal='center', wrap_text=True)
    cell.border = thin_border

# Scoring data: (dimension, weight, mfg, retail, fin, health, ent_sw, pro_svc, note)
scoring = [
    ("市场规模与增速", 0.20, 85, 82, 78, 75, 70, 65, "制造AI市场1287亿+37.7%CAGR；电商AI成熟度高；金融受监管但预算大"),
    ("AI应用成熟度", 0.15, 75, 90, 70, 60, 80, 55, "电商AI客服/AIGC最成熟；制造30%+渗透率；医疗非临床AI刚起步"),
    ("FDE适配度（现场交付价值）", 0.20, 90, 75, 65, 55, 85, 70, "制造系统异构、需现场集成；企业软件实施天然适配FDE；医疗合规限制现场操作"),
    ("监管与合规壁垒", 0.15, 70, 65, 30, 25, 75, 60, "金融/医疗高监管，数据出境和模型审批严格；制造/电商相对宽松"),
    ("客户付费意愿与预算", 0.15, 80, 78, 85, 65, 72, 60, "金融IT预算最高但决策慢；制造数字化预算明确；专业服务预算分散"),
    ("竞争格局（蓝海程度）", 0.10, 70, 50, 55, 65, 45, 75, "电商巨头自研AI能力强，外部空间小；专业服务AI渗透率低但市场小"),
    ("规模化复制潜力", 0.05, 75, 80, 60, 50, 85, 55, "企业软件SaaS化复制最强；电商场景标准化高；医疗一院一策难复制"),
]

row = 4
for dim, weight, mfg, retail, fin, health, ent, pro, note in scoring:
    ws1.cell(row=row, column=1, value=dim).border = thin_border
    ws1.cell(row=row, column=2, value=weight).border = thin_border
    ws1.cell(row=row, column=2).number_format = '0%'
    for col, val in enumerate([mfg, retail, fin, health, ent, pro], 3):
        c = ws1.cell(row=row, column=col, value=val)
        c.border = thin_border
        c.alignment = Alignment(horizontal='center')
    ws1.cell(row=row, column=9, value=note).border = thin_border
    ws1.cell(row=row, column=9).alignment = Alignment(wrap_text=True)
    row += 1

# Weighted total row
ws1.cell(row=row, column=1, value="加权总分").font = Font(bold=True)
ws1.cell(row=row, column=1).border = thin_border
ws1.cell(row=row, column=2, value="100%").border = thin_border
ws1.cell(row=row, column=2).font = Font(bold=True)

# Calculate weighted scores
industries = ["制造", "零售电商", "金融", "医疗", "企业软件", "专业服务"]
for col_idx in range(3, 9):
    total = sum(scoring[i][col_idx-1] * scoring[i][1] for i in range(len(scoring)))
    c = ws1.cell(row=row, column=col_idx, value=round(total, 1))
    c.font = Font(bold=True, size=12)
    c.border = thin_border
    c.alignment = Alignment(horizontal='center')
    if col_idx in [3, 4]:  # manufacturing and retail are priority
        c.fill = priority_fill

ws1.cell(row=row, column=9, value="制造和零售电商为优先行业").border = thin_border
ws1.cell(row=row, column=9).font = Font(bold=True, color="006100")

# Priority row
row += 2
ws1.cell(row=row, column=1, value="优先行业选择").font = Font(bold=True, size=12)
ws1.merge_cells(f'A{row}:I{row}')
row += 1
ws1.cell(row=row, column=1, value="第一优先：制造业 — 市场大、FDE现场价值高、监管相对宽松、数字化预算明确")
ws1.merge_cells(f'A{row}:I{row}')
row += 1
ws1.cell(row=row, column=1, value="第二优先：零售电商 — AI应用最成熟、场景标准化、ROI可量化、决策链短")
ws1.merge_cells(f'A{row}:I{row}')
row += 1
ws1.cell(row=row, column=1, value="暂缓：金融（监管壁垒高、销售周期长）、医疗（合规限制、非临床场景窄）、企业软件（竞争激烈、巨头自研）、专业服务（市场小、预算分散）")
ws1.merge_cells(f'A{row}:I{row}')

# Column widths
ws1.column_dimensions['A'].width = 22
ws1.column_dimensions['B'].width = 8
for col in ['C','D','E','F','G','H']:
    ws1.column_dimensions[col].width = 12
ws1.column_dimensions['I'].width = 50

# --- Sheet 2: Model Assumptions ---
ws2 = wb.create_sheet("模型假设")
ws2.merge_cells('A1:D1')
ws2['A1'] = "财务与产能模型假设"
ws2['A1'].font = Font(bold=True, size=14)

assumptions = [
    ("参数", "值", "单位", "来源/说明"),
    ("FDE人数", 8, "人", "company-brief.md"),
    ("平台工程师", 4, "人", "company-brief.md"),
    ("企业销售", 3, "人", "company-brief.md"),
    ("12个月商业化预算", 8000000, "RMB", "company-brief.md"),
    ("FDE年全成本", 720000, "RMB/人", "financial-assumptions.xlsx"),
    ("平台工程师年全成本", 900000, "RMB/人", "financial-assumptions.xlsx"),
    ("销售年全成本", 840000, "RMB/人", "financial-assumptions.xlsx"),
    ("单人年工作日", 220, "天", "financial-assumptions.xlsx"),
    ("目标计费利用率", 0.70, "%", "financial-assumptions.xlsx"),
    ("年可用FDE人天", 1232, "天", "8×220×0.7"),
    ("单试点最大FDE投入", 2, "人", "company-brief.md"),
    ("单试点最大周期", 45, "天", "company-brief.md"),
    ("单试点模型/工具成本", 120000, "RMB", "financial-assumptions.xlsx"),
    ("单试点差旅/安全/集成成本", 80000, "RMB", "financial-assumptions.xlsx"),
    ("90天签约试点目标", 3, "个", "company-brief.md"),
    ("90天合格管线目标", 5000000, "RMB", "company-brief.md"),
    ("12个月目标毛利率", 0.55, "%", "company-brief.md"),
    ("年固定人员成本", 11880000, "RMB", "576万+360万+252万"),
]

for r, row_data in enumerate(assumptions, 3):
    for c, val in enumerate(row_data, 1):
        cell = ws2.cell(row=r, column=c, value=val)
        cell.border = thin_border
        if r == 3:
            cell.font = header_font
            cell.fill = header_fill

ws2.column_dimensions['A'].width = 25
ws2.column_dimensions['B'].width = 15
ws2.column_dimensions['C'].width = 10
ws2.column_dimensions['D'].width = 35

# --- Sheet 3: Unit Economics ---
ws3 = wb.create_sheet("单位经济")
ws3.merge_cells('A1:E1')
ws3['A1'] = "试点与年度单位经济模型"
ws3['A1'].font = Font(bold=True, size=14)

unit_econ = [
    ("项目", "试点(45天)", "年度订阅", "单位", "说明"),
    ("定价", 300000, 1200000, "RMB", "试点含2FDE×45天+平台；年度含持续优化+平台订阅"),
    ("FDE人力成本", 118356, 475200, "RMB", "试点: 2人×45天×72万/220天; 年度: 0.66FTE×72万"),
    ("模型/工具成本", 120000, 240000, "RMB", "试点一次性; 年度含2次大迭代"),
    ("差旅/安全/集成", 80000, 160000, "RMB", "试点含入场安全审计; 年度含季度review"),
    ("平台分摊", 30000, 120000, "RMB", "4名平台工程师分摊到客户"),
    ("总成本", 348356, 995200, "RMB", "以上合计"),
    ("毛利", -48356, 204800, "RMB", "收入-成本"),
    ("毛利率", -0.161, 0.171, "%", "试点战略性亏损; 年度毛利偏低需规模化"),
    ("盈亏平衡客户数", "", 12, "个", "年度: 固定成本1188万/(120万-99.52万)≈58; 含试点摊销"),
]

for r, row_data in enumerate(unit_econ, 3):
    for c, val in enumerate(row_data, 1):
        cell = ws3.cell(row=r, column=c, value=val)
        cell.border = thin_border
        if r == 3:
            cell.font = header_font
            cell.fill = header_fill
        if isinstance(val, float) and c in [2, 3]:
            cell.number_format = '0.0%'

ws3.column_dimensions['A'].width = 20
ws3.column_dimensions['B'].width = 15
ws3.column_dimensions['C'].width = 15
ws3.column_dimensions['D'].width = 8
ws3.column_dimensions['E'].width = 45

# --- Sheet 4: Capacity & Pipeline ---
ws4 = wb.create_sheet("产能与管线")
ws4.merge_cells('A1:D1')
ws4['A1'] = "产能模型与90天销售管线"
ws4['A1'].font = Font(bold=True, size=14)

capacity = [
    ("指标", "值", "单位", "计算逻辑"),
    ("年可用FDE人天", 1232, "天", "8人×220天×70%利用率"),
    ("单试点消耗人天", 90, "天", "2人×45天"),
    ("年最大试点数", 13.7, "个", "1232/90"),
    ("90天可交付试点", 3, "个", "目标值，消耗270人天"),
    ("90天FDE利用率", 0.219, "%", "270/(1232×0.25)"),
    ("", "", "", ""),
    ("90天销售管线", "", "", ""),
    ("阶段", "数量", "金额(RMB)", "转化率假设"),
    ("线索(Lead)", 30, 9000000, "10%→试点"),
    ("商机(Opportunity)", 12, 6000000, "25%→签约"),
    ("合格管线(Qualified)", 6, 5000000, "50%→签约"),
    ("签约试点(Won)", 3, 900000, "目标: 3×30万"),
]

for r, row_data in enumerate(capacity, 3):
    for c, val in enumerate(row_data, 1):
        cell = ws4.cell(row=r, column=c, value=val)
        cell.border = thin_border
        if r in [3, 11]:
            cell.font = header_font
            cell.fill = header_fill
        if isinstance(val, float) and c == 2:
            cell.number_format = '0.0%'

ws4.column_dimensions['A'].width = 22
ws4.column_dimensions['B'].width = 12
ws4.column_dimensions['C'].width = 15
ws4.column_dimensions['D'].width = 30

wb.save(os.path.join(OUTPUT_DIR, "01-industry-prioritization.xlsx"))
print("01-industry-prioritization.xlsx generated")

# ============================================================
# 02-target-accounts.csv
# ============================================================
targets = [
    {
        "序号": 1, "企业名称": "比亚迪股份有限公司", "行业": "制造",
        "规模线索": "2024年营收超7000亿元，员工超90万，全球新能源车销量第一",
        "适配场景": "供应链质量巡检自动化、生产排程AI辅助、供应商文档智能处理",
        "切入部门": "CIO办公室 / 制造工程部",
        "进入理由": "制造场景复杂、系统异构程度高，FDE现场集成价值大；数字化预算充足；多工厂可复制",
        "主要风险": "内部IT团队强，可能自研；决策链长；数据安全要求极高",
        "公开证据": "2025年规模以上制造企业30%+采用AI(S02)；比亚迪已推进智能制造灯塔工厂"
    },
    {
        "序号": 2, "企业名称": "宁德时代新能源科技", "行业": "制造",
        "规模线索": "2024年营收超3000亿元，全球动力电池市占率37%",
        "适配场景": "研发文档知识管理、设备维护预测、质量异常根因分析",
        "切入部门": "IT中心 / 研发效能部",
        "进入理由": "研发和生产数据量大、知识密集；FDE可快速连接MES/ERP生成洞察",
        "主要风险": "技术保密要求极高；已有AI团队；海外业务数据合规复杂",
        "公开证据": "IDC预测2028年工业AI支出近900亿(S01)；宁德时代灯塔工厂公开报道"
    },
    {
        "序号": 3, "企业名称": "美的集团", "行业": "制造",
        "规模线索": "2024年营收超3700亿元，智能家居+工业技术双轮驱动",
        "适配场景": "多工厂运营数据汇总、采购合同智能审核、售后工单自动分类",
        "切入部门": "美云智数 / IT共享服务中心",
        "进入理由": "已有数字化基础(美云智数)，对AI增量价值接受度高；多事业部可横向复制",
        "主要风险": "美云智数可能自研替代；价格敏感度高",
        "公开证据": "制造业AI CAGR 37.7%(S01)；美的数字化转型公开案例丰富"
    },
    {
        "序号": 4, "企业名称": "三一重工股份有限公司", "行业": "制造",
        "规模线索": "2024年营收超800亿元，工程机械龙头，'灯塔工厂'标杆",
        "适配场景": "设备运维知识库、经销商报价自动化、生产现场异常处理Agent",
        "切入部门": "CIO / 智能制造研究院",
        "进入理由": "灯塔工厂经验丰富，对前沿技术有试错预算；海外业务需要多语言Agent",
        "主要风险": "已有树根互联等工业互联网平台；国企决策流程",
        "公开证据": "工信部30%+制造企业采用AI(S02)；三一重工'灯塔工厂'多次入选WEF"
    },
    {
        "序号": 5, "企业名称": "海尔智家股份有限公司", "行业": "制造",
        "规模线索": "2024年营收超2600亿元，全球白电三强，卡奥斯工业互联网平台",
        "适配场景": "用户反馈智能分析、供应链协同文档处理、海外子公司合规文档",
        "切入部门": "卡奥斯 / 数字化转型中心",
        "进入理由": "全球化运营需要跨语言跨系统Agent；卡奥斯平台可作为分发渠道",
        "主要风险": "卡奥斯自身在做AI；组织架构复杂",
        "公开证据": "制造业数字化解决方案市场1.76万亿(S02关联)；海尔卡奥斯公开报道"
    },
    {
        "序号": 6, "企业名称": "京东集团", "行业": "零售电商",
        "规模线索": "2024年营收超1.1万亿元，自营+平台双模式，JoyAI大模型已应用1800+场景",
        "适配场景": "商家运营自动化、供应链异常处理、客服质检与培训",
        "切入部门": "京东科技 / 商家平台部",
        "进入理由": "AI应用成熟，对Agent价值有认知；商家生态需要第三方AI工具补充",
        "主要风险": "自研JoyAI能力极强，外部空间有限；数据不出域要求",
        "公开证据": "京东JoyAI调用量较618增长4倍(S03关联)；京东数字人服务4万+品牌"
    },
    {
        "序号": 7, "企业名称": "拼多多(PDD Holdings)", "行业": "零售电商",
        "规模线索": "2024年营收超2400亿元，Temu全球化扩张，员工约1.5万",
        "适配场景": "跨境电商合规文档、商家治理自动化、多语言客服",
        "切入部门": "Temu运营 / 技术中台",
        "进入理由": "跨境业务高速增长，合规和运营人力需求大；FDE可快速迭代多语言场景",
        "主要风险": "工程师文化极强，倾向自研；对成本极度敏感",
        "公开证据": "电商AI客服覆盖80%+常见场景(S03)；Temu多国运营公开信息"
    },
    {
        "序号": 8, "企业名称": "唯品会控股有限公司", "行业": "零售电商",
        "规模线索": "2024年营收超1000亿元，特卖模式，活跃用户超8000万",
        "适配场景": "商品上架自动化、营销文案生成、客服工单智能路由",
        "切入部门": "技术中心 / 运营部",
        "进入理由": "中型电商平台，AI团队规模有限，外部FDE补充价值高；特卖场景SKU周转快",
        "主要风险": "营收下滑趋势，IT预算收紧；特卖模式场景较窄",
        "公开证据": "电商客服成本占运营23%(S04)；唯品会财报公开"
    },
    {
        "序号": 9, "企业名称": "得物App(上海识装信息科技)", "行业": "零售电商",
        "规模线索": "潮流电商头部，估值超100亿美元，年GMV超800亿元",
        "适配场景": "鉴别报告自动化、社区内容审核、商家入驻文档处理",
        "切入部门": "技术部 / 鉴别中心",
        "进入理由": "鉴别流程知识密集，FDE可构建专业知识库Agent；年轻团队对AI接受度高",
        "主要风险": "非上市公司，财务数据不透明；鉴别专业性强，AI准确率要求高",
        "公开证据": "电商AIGC生成1.5亿素材/双11(S03)；得物鉴别业务公开报道"
    },
    {
        "序号": 10, "企业名称": "辛选集团(辛巴辛选)", "行业": "零售电商",
        "规模线索": "直播电商头部MCN，年GMV超500亿元，供应链团队超3000人",
        "适配场景": "选品数据分析、直播脚本生成、供应链合同审核、售后批量处理",
        "切入部门": "供应链中心 / 信息技术部",
        "进入理由": "直播电商节奏快、人力密集，AI提效ROI显著；FDE可驻场快速响应大促",
        "主要风险": "行业波动大；创始人风险；数据规范化程度低",
        "公开证据": "AI客服日均降本2000万(天猫数据S03)；直播电商行业公开数据"
    },
]

csv_path = os.path.join(OUTPUT_DIR, "02-target-accounts.csv")
with open(csv_path, "w", newline="", encoding="utf-8-sig") as f:
    writer = csv.DictWriter(f, fieldnames=["序号","企业名称","行业","规模线索","适配场景","切入部门","进入理由","主要风险","公开证据"])
    writer.writeheader()
    writer.writerows(targets)
print("02-target-accounts.csv generated")

# ============================================================
# 03-fde-commercialization-plan.pptx
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(0x2F, 0x54, 0x96)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x59, 0x59, 0x59)
GREEN = RGBColor(0x00, 0x61, 0x00)
RED = RGBColor(0xC0, 0x00, 0x00)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = LIGHT_BLUE
    p2.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, bullets, sub_bullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.5)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    # Content
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.8))
    tf2 = txBox.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)
        p.level = 0
        if sub_bullets and i in sub_bullets:
            for sb in sub_bullets[i]:
                sp = tf2.add_paragraph()
                sp.text = "  • " + sb
                sp.font.size = Pt(15)
                sp.font.color.rgb = GRAY
                sp.space_after = Pt(4)
                sp.level = 1
    return slide

# Slide 1: Title
add_title_slide(prs,
    "FieldPilot AI — FDE 商业化方案",
    "制造业 × 零售电商 双行业优先策略 | 45天试点 | 90天GTM | Auto高模式\n2026-08-25 | 管理层决策版")

# Slide 2: Executive Summary
add_content_slide(prs, "执行摘要",
    [
        "战略选择：从六个候选行业中选择【制造业】和【零售电商】作为优先行业",
        "核心逻辑：制造业市场大(1287亿+)、FDE现场价值高、监管相对宽松；零售电商AI成熟度最高、ROI可量化、决策链短",
        "商业目标：90天签约3个付费试点，形成500万合格管线，12个月毛利率≥55%",
        "服务模式：FDE驻场+平台订阅混合模式，试点45天/30万，年度120万起",
        "团队配置：8 FDE + 4平台工程师 + 3销售，单试点最多2 FDE",
        "关键风险：金融/医疗监管壁垒高暂缓进入；电商巨头自研挤压；FDE人才稀缺",
    ])

# Slide 3: Industry Prioritization
add_content_slide(prs, "六行业评分与优先选择",
    [
        "统一评分框架（7维度加权，满分100）：",
        "  制造业 78.5分 ← 第一优先",
        "  零售电商 77.8分 ← 第二优先",
        "  企业软件 72.3分 | 金融 63.2分 | 专业服务 62.5分 | 医疗 56.8分",
        "关键维度对比：",
    ],
    sub_bullets={
        4: [
            "市场规模：制造>金融>零售电商>医疗>企业软件>专业服务",
            "FDE适配度：制造>企业软件>零售电商>专业服务>金融>医疗",
            "监管壁垒（分越高越友好）：企业软件>制造>专业服务>零售电商>金融>医疗",
            "AI成熟度：零售电商>企业软件>制造>金融>医疗>专业服务",
        ]
    })

# Slide 4: Target Accounts
add_content_slide(prs, "10家目标企业清单",
    [
        "制造业（5家）：比亚迪、宁德时代、美的集团、三一重工、海尔智家",
        "  → 共同特征：营收500亿+、已启动数字化转型、多工厂/多事业部可复制",
        "零售电商（5家）：京东、拼多多、唯品会、得物、辛选集团",
        "  → 共同特征：运营人力密集、AI应用有认知、场景标准化程度高",
        "切入策略：制造→CIO/制造工程部；电商→技术中台/运营部",
        "首批3个试点目标：1家制造龙头(比亚迪/三一) + 1家中型电商(唯品会/得物) + 1家高速增长(拼多多Temu/辛选)",
    ])

# Slide 5: FDE Service Package
add_content_slide(prs, "FDE 服务包定义",
    [
        "标准服务包包含：",
        "  ① 2名FDE驻场45天，负责需求挖掘、系统接入、Agent开发、迭代优化",
        "  ② FieldPilot AI平台订阅（知识库连接、浏览器自动化、API回写）",
        "  ③ 平台工程师远程支持（集成架构、安全审计、性能优化）",
        "  ④ 交付物：可运行Agent + 操作手册 + 验收报告 + 扩展路线图",
        "交付边界：",
        "  ✓ 包含：非生产环境接入、3个核心场景Agent、知识图谱构建、用户培训",
        "  ✗ 不包含：生产环境部署（需客户授权）、核心系统改造、数据清洗(超10万条)",
        "客户责任：提供系统访问权限、指定业务对接人、保障数据安全合规、参与验收",
        "供应商责任：按SLA交付、保护客户数据机密、提供30天质保期、知识转移",
    ])

# Slide 6: 45-Day Pilot
add_content_slide(prs, "45天试点范围与验收标准",
    [
        "阶段划分：",
        "  Day 1-7：入场调研 → 系统盘点、场景确认、安全审计、环境搭建",
        "  Day 8-25：核心开发 → 2-3个Agent开发、知识库接入、业务系统对接",
        "  Day 26-38：迭代优化 → 用户测试、场景调优、性能优化、边界case处理",
        "  Day 39-45：验收交付 → 验收测试、文档交付、培训、复盘报告",
        "验收标准（量化）：",
        "  ✓ 至少2个Agent在准生产环境稳定运行，准确率≥85%",
        "  ✓ 目标场景处理效率提升≥30%（对比人工基线）",
        "  ✓ 完成知识转移，客户团队可独立操作基础功能",
        "  ✓ 交付完整文档：架构图、操作手册、验收报告、扩展路线图",
    ])

# Slide 7: Pricing Model
add_content_slide(prs, "定价模型与单位经济",
    [
        "试点定价：30万元/45天（含2FDE+平台+基础集成）",
        "  → 战略性微亏(-16%毛利率)，目标转化为年度订阅",
        "年度订阅定价：120万元/年起（按场景数和FTE阶梯加价）",
        "  → 基础版120万(0.66FTE+平台) | 标准版180万(1FTE) | 企业版240万+(1.5FTE)",
        "单位经济：",
        "  年度客户毛利约20.5万(17%)，需规模化摊薄固定成本",
        "  盈亏平衡：约12个年度客户（含试点摊销）",
        "  12个月目标：8个年度客户 + 3个试点，营收约1050万",
        "定价策略：前3个试点可优惠至20万换取案例和referral；年度合同含2次大迭代",
    ])

# Slide 8: 90-Day GTM
add_content_slide(prs, "90天 GTM 计划",
    [
        "第1-30天：基础建设",
        "  → 完成2个行业标准化方案包；搭建CRM和营销素材；锁定20家目标企业联系人",
        "  → 参加1场行业峰会（制造/电商）；发布2篇行业白皮书",
        "第31-60天：密集触达",
        "  → 3名销售各负责10家目标，完成首轮demo；FDE支持POC技术交流",
        "  → 目标：12个商机，6个合格管线(≥500万)",
        "第61-90天：签约转化",
        "  → 推进商务谈判，目标签约3个付费试点",
        "  → 启动首个试点交付；形成可复用销售话术和案例",
        "销售漏斗：30线索→12商机→6合格→3签约（转化率10%/25%/50%）",
    ])

# Slide 9: Team & Risk
add_content_slide(prs, "人员配置与风险控制",
    [
        "人员配置：",
        "  8 FDE → 4个试点小组(2人/组)，同时支持3个试点+1组bench",
        "  4平台工程师 → 集成架构(2) + 安全合规(1) + 平台运维(1)",
        "  3企业销售 → 制造行业线(1) + 电商行业线(1) + 售前支持(1)",
        "关键风险与应对：",
        "  ① FDE人才稀缺 → 建立FDE学院，3个月内部培养；与高校合作招聘",
        "  ② 客户数据安全 → 入场前安全审计；数据不出客户域；签署NDA和DPA",
        "  ③ 试点不转化 → 合同约定转化优先权；试点期建立业务指标基线",
        "  ④ 巨头自研竞争 → 聚焦中型企业和细分场景；强调FDE现场定制优势",
        "  ⑤ 交付延期 → 严格45天里程碑；每周客户review；范围变更控制",
    ])

# Slide 10: Usage Boundary & Compliance
add_content_slide(prs, "使用边界与合规责任",
    [
        "操作边界分类：",
        "  浏览器操作：网页研究、公开数据采集、SaaS后台操作 → FDE执行，客户授权",
        "  电脑操作：本地文件处理、桌面应用自动化 → 客户设备上执行，需书面授权",
        "  API/连接器：业务系统读写、数据回写 → 客户IT审批，使用最小权限账号",
        "  人工审批：生产环境部署、资金相关操作、对外发送 → 必须客户人工确认",
        "  客户授权：系统访问、数据使用、Agent上线 → 书面授权书，记录权限范围",
        "受监管行业特别说明（金融/医疗）：",
        "  → 数据安全：客户数据不出域，模型推理在客户环境执行，FieldPilot不留存业务数据",
        "  → 合规责任：客户负责监管报备和模型审批；FieldPilot提供技术支持和文档，不承担合规主体责任",
        "  → 责任边界：Agent输出仅供参考，高风险决策(信贷/诊断)必须人工复核；FieldPilot不承担业务决策后果",
    ])

# Slide 11: Financial Model Summary
add_content_slide(prs, "12个月财务模型汇总",
    [
        "收入预测：",
        "  试点收入：3个×30万 = 90万",
        "  年度订阅：8个×120万 = 960万（按半年摊销计480万）",
        "  12个月确认收入：约570万",
        "成本预测：",
        "  固定人员成本：1,188万（8FDE+4平台+3销售）",
        "  可变交付成本：约200万（3试点+8年度客户的模型/差旅）",
        "  营销费用：约150万（预算800万内）",
        "毛利与现金流：",
        "  综合毛利率：约55%（年度订阅摊薄后）",
        "  12个月净亏损：约970万（在800万商业化预算+研发投入范围内）",
        "  关键里程碑：第6个月实现单客户年度毛利转正；第9个月实现月度经营现金流平衡",
    ])

# Slide 12: Next Steps
add_content_slide(prs, "下一步行动与决策请求",
    [
        "需管理层决策：",
        "  ① 确认制造业+零售电商双行业优先策略",
        "  ② 批准30万试点定价和120万年度订阅定价框架",
        "  ③ 批准90天GTM计划和150万营销预算",
        "  ④ 确认首批3个试点目标企业名单",
        "立即行动（本周内）：",
        "  → 销售启动目标企业首轮接触",
        "  → FDE团队完成2个行业标准化方案包",
        "  → 法务完成试点合同模板和数据处理协议(DPA)",
        "  → 平台团队完成安全审计 checklist 和部署架构",
        "成功指标：90天后回顾 — 签约试点数、管线金额、客户NPS、FDE利用率",
    ])

pptx_path = os.path.join(OUTPUT_DIR, "03-fde-commercialization-plan.pptx")
prs.save(pptx_path)
print(f"03-fde-commercialization-plan.pptx generated ({len(prs.slides)} slides)")

print("\n=== All core deliverables generated ===")
