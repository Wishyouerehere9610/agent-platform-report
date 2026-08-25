#!/usr/bin/env python3
"""Generate 03-fde-commercialization-plan.pptx for FieldPilot AI."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
from pathlib import Path

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT_BLUE = RGBColor(0x2F, 0x54, 0x96)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
RED = RGBColor(0xE0, 0x53, 0x3E)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_box(slide, text, left=Inches(0.6), top=Inches(0.3), width=Inches(12), height=Inches(0.8)):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = 'Microsoft YaHei'
    return shape

def add_subtitle(slide, text, top=Inches(1.1)):
    shape = slide.shapes.add_textbox(Inches(0.6), top, Inches(12), Inches(0.5))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = 'Microsoft YaHei'
    return shape

def add_text_box(slide, text, left, top, width, height, font_size=12, color=DARK_GRAY, bold=False):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Microsoft YaHei'
    return shape

def add_shape_box(slide, left, top, width, height, fill_color, text="", text_color=WHITE, font_size=11, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = 'Microsoft YaHei'
    return shape

def add_table(slide, rows_data, left, top, width, height, header_color=ACCENT_BLUE, col_widths=None):
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    
    for r, row_data in enumerate(rows_data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9) if r > 0 else Pt(10)
                p.font.name = 'Microsoft YaHei'
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = DARK_GRAY
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    return table_shape

# ============ Slide 1: Title ============
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1, DARK_BLUE)

shape = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "FieldPilot AI"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = 'Microsoft YaHei'
p.alignment = PP_ALIGN.CENTER

shape2 = slide1.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
tf2 = shape2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "FDE 商业化方案"
p2.font.size = Pt(32)
p2.font.color.rgb = RGBColor(0x9D, 0xC3, 0xE6)
p2.font.name = 'Microsoft YaHei'
p2.alignment = PP_ALIGN.CENTER

shape3 = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
tf3 = shape3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "面向管理层决策的行业优先级、目标客户、定价与GTM方案"
p3.font.size = Pt(16)
p3.font.color.rgb = RGBColor(0xB4, 0xC7, 0xE7)
p3.font.name = 'Microsoft YaHei'
p3.alignment = PP_ALIGN.CENTER

shape4 = slide1.shapes.add_textbox(Inches(1), Inches(6), Inches(11), Inches(0.8))
tf4 = shape4.text_frame
p4 = tf4.paragraphs[0]
p4.text = "2026年8月  |  Forward Deployed Engineer 商业化"
p4.font.size = Pt(12)
p4.font.color.rgb = RGBColor(0x8D, 0xA0, 0xB5)
p4.font.name = 'Microsoft YaHei'
p4.alignment = PP_ALIGN.CENTER

# ============ Slide 2: Executive Summary ============
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)
add_title_box(slide2, "执行摘要")
add_subtitle(slide2, "核心结论与关键数字")

# Key metrics boxes
metrics = [
    ("2", "优先行业", GREEN),
    ("10", "目标企业", ACCENT_BLUE),
    ("150万", "年度定价/客户", ORANGE),
    ("78%", "综合毛利率", GREEN),
]
for i, (num, label, color) in enumerate(metrics):
    left = Inches(0.6 + i * 3.1)
    add_shape_box(slide2, left, Inches(1.6), Inches(2.8), Inches(1.5), color, num, WHITE, 28, True)
    add_text_box(slide2, label, left, Inches(3.2), Inches(2.8), Inches(0.4), 12, DARK_GRAY, True)

# Summary text
summary_text = (
    "行业选择：基于7维加权评分框架，从6个候选行业中选定制造业（83.0分）和零售电商（84.7分）为优先行业。\n"
    "商业目标：90天内签约3个付费试点（450万合同额+500万管线），12个月目标综合毛利率≥55%。\n"
    "定价模型：试点35-50万/45天，年度60-180万/年，含FDE服务包+平台订阅。\n"
    "团队配置：8名FDE年可用1232人天，支持13.7个试点等效，90天目标消耗270人天。\n"
    "合规边界：受监管行业（金融、医疗）暂缓，优先行业以非核心决策流程切入，所有AI输出经人工审核。"
)
add_text_box(slide2, summary_text, Inches(0.6), Inches(3.8), Inches(12), Inches(3.2), 13, DARK_GRAY)

# ============ Slide 3: Industry Prioritization ============
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3)
add_title_box(slide3, "行业优先级评分")
add_subtitle(slide3, "7维加权评分框架 | 6个候选行业 | 选定2个优先行业")

table_data = [
    ["行业", "市场规模(20%)", "AI适配(20%)", "交付可行(15%)", "监管(15%)", "预算(10%)", "竞争(10%)", "产能(10%)", "总分", "结论"],
    ["制造业", "18", "17", "14", "12", "9", "7", "9", "83.0", "优先"],
    ["零售电商", "17", "19", "14", "13", "8", "6", "9", "84.7", "优先"],
    ["金融服务", "16", "14", "10", "6", "9", "7", "7", "71.5", "暂缓"],
    ["医疗健康", "14", "13", "9", "5", "7", "6", "6", "63.0", "暂缓"],
    ["企业软件", "15", "15", "12", "11", "7", "4", "8", "74.0", "备选"],
    ["专业服务", "10", "16", "13", "10", "6", "8", "7", "71.0", "备选"],
]
add_table(slide3, table_data, Inches(0.4), Inches(1.6), Inches(12.5), Inches(3.5),
          col_widths=[Inches(1.2)] + [Inches(1.25)]*8 + [Inches(0.8)])

add_text_box(slide3, 
    "评分依据：市场规模（工信部/艾瑞/IIM报告）| AI适配度（FieldPilot能力匹配分析）| 交付可行性（2 FDE/45天约束）\n"
    "监管复杂度（金融监管总局/国家药监局/网信办公开文件）| 客户预算（行业IT支出公开数据）| 竞争烈度（市场观察）| 产能匹配（1232人天/年）",
    Inches(0.4), Inches(5.3), Inches(12.5), Inches(1.5), 10, RGBColor(0x80, 0x80, 0x80))

# ============ Slide 4: Priority Industry Rationale ============
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4)
add_title_box(slide4, "优先行业选择理由")
add_subtitle(slide4, "制造业 + 零售电商 | 基于市场规模、AI适配度与交付可行性")

# Two columns
add_shape_box(slide4, Inches(0.5), Inches(1.6), Inches(6), Inches(0.6), GREEN, "制造业（加权 83.0 分）", WHITE, 16)
mfg_text = (
    "市场规模：数字化转型1.76万亿(2025)→2.01万亿(2026)，增速14%+\n"
    "AI适配：采购比价、质量文档、供应链分析与FieldPilot高度匹配\n"
    "政策红利：AI+制造专项实施意见，智能工厂梯度培育\n"
    "买家：COO/CIO/工厂运营，预算明确\n"
    "监管：中等复杂度，数据隔离要求可控\n"
    "切入点：中型制造企业（年收入10-50亿）\n"
    "场景：采购自动化→质量报告→供应链研究\n"
    "来源：工信部、艾瑞咨询、中宏网（4个来源）"
)
add_text_box(slide4, mfg_text, Inches(0.5), Inches(2.3), Inches(6), Inches(4.5), 11, DARK_GRAY)

add_shape_box(slide4, Inches(6.8), Inches(1.6), Inches(6), Inches(0.6), ACCENT_BLUE, "零售电商（加权 84.7 分）", WHITE, 16)
retail_text = (
    "市场规模：智慧零售300亿(2026)，即时零售万亿级\n"
    "AI适配：选品分析、内容生成、客服自动化高度匹配\n"
    "AI渗透率：超80%企业已部署AI，45.3%高频使用\n"
    "买家：电商总经理/CMO/COO，决策快\n"
    "监管：消费者数据保护与广告合规，门槛中等\n"
    "切入点：中大型连锁零售和品牌电商\n"
    "场景：内容生成→选品分析→客服自动化\n"
    "来源：DoNews、IIM、CCFA、中商产业研究院（4个来源）"
)
add_text_box(slide4, retail_text, Inches(6.8), Inches(2.3), Inches(6), Inches(4.5), 11, DARK_GRAY)

# ============ Slide 5: Target Accounts ============
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5)
add_title_box(slide5, "目标企业清单")
add_subtitle(slide5, "10家企业 | 5家制造业 + 5家零售电商 | 基于公开年报与行业案例")

target_data = [
    ["#", "企业", "行业", "规模线索", "适配场景", "切入部门"],
    ["1", "美的集团", "制造", "3000亿+", "采购比价/质量文档/供应链分析", "COO/CIO"],
    ["2", "三一重工", "制造", "800亿+", "设备维护知识库/供应链/采购", "CIO/灯塔工厂"],
    ["3", "宁德时代", "制造", "3600亿+", "供应链研究/质量文档/合规报告", "CIO/供应链"],
    ["4", "海尔智家", "制造", "2000亿+", "供应链/内容生成/客服自动化", "CMO/CIO"],
    ["5", "汇川技术", "制造", "300亿+", "采购比价/技术文档/供应链", "CIO/运营"],
    ["6", "永辉超市", "零售", "700亿+", "选品分析/库存分析/客服", "COO/CIO"],
    ["7", "孩子王", "零售", "100亿+", "选品/内容生成/会员运营", "CMO/CIO"],
    ["8", "百果园", "零售", "100亿+", "供应链/选品/门店运营报告", "COO/CIO"],
    ["9", "锅圈食汇", "零售", "50亿+", "选品/供应链/营销内容", "COO/CIO"],
    ["10", "银泰商业", "零售", "100亿+", "选品/内容生成/客服/竞品", "CMO/CIO"],
]
add_table(slide5, target_data, Inches(0.3), Inches(1.6), Inches(12.7), Inches(5),
          col_widths=[Inches(0.5), Inches(1.5), Inches(0.8), Inches(1.2), Inches(5.5), Inches(1.5)])

add_text_box(slide5, "注：企业规模基于公开年报，不构成采购意向或合作状态的断言。详见 02-target-accounts.csv。",
    Inches(0.3), Inches(6.8), Inches(12), Inches(0.5), 9, RGBColor(0x80, 0x80, 0x80))

# ============ Slide 6: FDE Service Package & Pilot Scope ============
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6)
add_title_box(slide6, "FDE 服务包与试点范围")
add_subtitle(slide6, "45天交付周期 | 2名FDE/试点 | 明确验收标准与责任边界")

add_shape_box(slide6, Inches(0.5), Inches(1.6), Inches(3.8), Inches(0.5), ACCENT_BLUE, "服务包内容", WHITE, 13)
svc_text = (
    "1. 工作流诊断与方案设计（第1-5天）\n"
    "2. 系统接入与数据连接器配置（第6-15天）\n"
    "3. AI Agent 定制与场景验证（第16-30天）\n"
    "4. 用户培训与流程文档（第31-40天）\n"
    "5. 验收交付与改进建议（第41-45天）"
)
add_text_box(slide6, svc_text, Inches(0.5), Inches(2.2), Inches(3.8), Inches(3), 10, DARK_GRAY)

add_shape_box(slide6, Inches(4.5), Inches(1.6), Inches(4), Inches(0.5), GREEN, "验收标准", WHITE, 13)
acc_text = (
    "1. ≥3个核心场景AI Agent稳定运行\n"
    "2. 单场景自动化率≥60%\n"
    "3. 端到端流程文档完整交付\n"
    "4. 客户方≥2人完成培训认证\n"
    "5. 安全审计通过（数据隔离验证）\n"
    "6. 客户满意度评分≥4.0/5.0"
)
add_text_box(slide6, acc_text, Inches(4.5), Inches(2.2), Inches(4), Inches(3), 10, DARK_GRAY)

add_shape_box(slide6, Inches(8.7), Inches(1.6), Inches(4), Inches(0.5), ORANGE, "责任边界", WHITE, 13)
resp_text = (
    "FieldPilot 负责：\n"
    "• Agent 配置与连接器开发\n"
    "• 场景验证与流程文档\n"
    "• 安全合规技术方案\n\n"
    "客户 负责：\n"
    "• 系统访问授权与数据提供\n"
    "• 业务专家参与场景定义\n"
    "• 验收测试与签字确认"
)
add_text_box(slide6, resp_text, Inches(8.7), Inches(2.2), Inches(4), Inches(4), 10, DARK_GRAY)

# ============ Slide 7: Pricing & Financial Model ============
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7)
add_title_box(slide7, "定价与财务模型")
add_subtitle(slide7, "试点定价 + 年度订阅 | 收入、成本、产能与毛利模型")

price_data = [
    ["项目", "试点(万/45天)", "年度(万/年)", "成本(万)", "毛利率", "说明"],
    ["FDE服务包-标准", "35", "120", "20", "43%", "2 FDE×45天"],
    ["FDE服务包-高级", "50", "180", "25", "50%", "含平台工程师支持"],
    ["平台订阅-基础", "—", "60", "10", "83%", "基础连接器+SSO"],
    ["平台订阅-企业", "—", "120", "15", "88%", "高级连接器+审计"],
    ["加权综合", "40", "150", "30", "80%", "目标≥55%毛利率"],
]
add_table(slide7, price_data, Inches(0.5), Inches(1.6), Inches(12), Inches(2.5),
          col_widths=[Inches(2.5), Inches(2), Inches(2), Inches(1.5), Inches(1.5), Inches(2.5)])

# Revenue projection
add_shape_box(slide7, Inches(0.5), Inches(4.3), Inches(5.8), Inches(0.5), ACCENT_BLUE, "12个月收入预测", WHITE, 13)
rev_text = (
    "0-90天：3个试点×150万/年 = 450万合同额 + 500万管线\n"
    "90-180天：3个试点 = 450万 + 600万管线\n"
    "180-365天：4个试点 = 600万 + 800万管线\n"
    "12个月合计：10个试点，1500万合同额，1900万管线\n"
    "消耗FDE人天：900天（占可用1232天的73%）"
)
add_text_box(slide7, rev_text, Inches(0.5), Inches(4.9), Inches(5.8), Inches(2.2), 10, DARK_GRAY)

add_shape_box(slide7, Inches(6.5), Inches(4.3), Inches(6), Inches(0.5), GREEN, "产能与毛利", WHITE, 13)
cap_text = (
    "年可用FDE人天：1,232天（8人×220天×70%）\n"
    "年最大试点等效：13.7个（1,232÷90人天/试点）\n"
    "年固定人员成本：1,188万（FDE 576+平台360+销售252）\n"
    "90天可变交付成本：60万（3试点×20万）\n"
    "综合毛利率：约78%（含平台订阅高毛利部分）\n"
    "盈亏平衡：约8个年度合同（8×150万=1,200万≈固定成本）"
)
add_text_box(slide7, cap_text, Inches(6.5), Inches(4.9), Inches(6), Inches(2.2), 10, DARK_GRAY)

# ============ Slide 8: 90-Day GTM ============
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8)
add_title_box(slide8, "90天 GTM 计划")
add_subtitle(slide8, "销售漏斗 | 实施节奏 | 人员配置 | 风险控制")

# Timeline
phases = [
    ("第1-15天", "线索开发", ACCENT_BLUE, "销售3人各负责1行业\n目标30个有效线索\n10家进入初步接触\n3家进入方案演示"),
    ("第16-30天", "方案演示", GREEN, "FDE 2人参与\n5场方案演示\n3家进入试点谈判\n签约1家付费试点"),
    ("第31-60天", "试点交付", ORANGE, "2 FDE驻场交付\n第1个试点启动\n签约第2个试点\n完成1个场景验证"),
    ("第61-90天", "验收与扩展", RED, "第1个试点验收\n签约第3个试点\n形成可复制案例\n500万有效管线"),
]
for i, (phase, title, color, desc) in enumerate(phases):
    left = Inches(0.4 + i * 3.2)
    add_shape_box(slide8, left, Inches(1.6), Inches(0.8), Inches(0.5), color, phase, WHITE, 10)
    add_shape_box(slide8, left, Inches(2.2), Inches(3), Inches(0.5), color, title, WHITE, 13)
    add_text_box(slide8, desc, left, Inches(2.8), Inches(3), Inches(2), 10, DARK_GRAY)

# Sales funnel
add_shape_box(slide8, Inches(0.5), Inches(5), Inches(5.8), Inches(0.5), ACCENT_BLUE, "销售漏斗", WHITE, 13)
funnel_text = (
    "30个有效线索 → 15家初步接触 → 8家方案演示 → 5家试点谈判 → 3个签约试点\n"
    "转化率：线索→接触 50% | 接触→演示 53% | 演示→谈判 63% | 谈判→签约 60%\n"
    "平均销售周期：60天（从首次接触到签约）"
)
add_text_box(slide8, funnel_text, Inches(0.5), Inches(5.6), Inches(5.8), Inches(1.5), 10, DARK_GRAY)

add_shape_box(slide8, Inches(6.5), Inches(5), Inches(6), Inches(0.5), RED, "风险控制", WHITE, 13)
risk_text = (
    "1. 试点前签署数据安全协议与责任边界确认\n"
    "2. FDE驻场期间日报+周报机制，客户可随时中止\n"
    "3. 试点范围限定非核心决策流程，AI输出经人工审核\n"
    "4. 单试点超45天或超2 FDE须管理层审批\n"
    "5. 合规红线：不处理受监管数据、不代客决策"
)
add_text_box(slide8, risk_text, Inches(6.5), Inches(5.6), Inches(6), Inches(1.5), 10, DARK_GRAY)

# ============ Slide 9: Implementation & Staffing ============
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9)
add_title_box(slide9, "实施节奏与人员配置")
add_subtitle(slide9, "8名FDE | 4名平台工程师 | 3名企业销售 | 90天节奏")

staff_data = [
    ["角色", "人数", "90天任务分配", "关键交付物"],
    ["FDE", "8", "2人×3试点=6人驻场交付 | 2人做售前支持", "3个试点验收报告+流程文档"],
    ["平台工程师", "4", "1人连接器开发 | 1人安全合规 | 2人平台运维", "连接器组件+安全审计报告"],
    ["企业销售", "3", "2人制造业 | 1人零售电商 | 各负责15个线索", "3个签约合同+500万管线"],
]
add_table(slide9, staff_data, Inches(0.5), Inches(1.6), Inches(12), Inches(2.5),
          col_widths=[Inches(2), Inches(1), Inches(5.5), Inches(3.5)])

add_shape_box(slide9, Inches(0.5), Inches(4.3), Inches(5.8), Inches(0.5), ACCENT_BLUE, "实施节奏", WHITE, 13)
rhythm_text = (
    "第1周：客户启动会、系统访问授权、数据安全协议\n"
    "第2-3周：工作流诊断、连接器配置、Agent定制\n"
    "第4-6周：场景验证、用户培训、流程文档编写\n"
    "第7周：验收测试、安全审计、客户签字确认\n"
    "第8周起：交付总结、案例包装、复制方案输出"
)
add_text_box(slide9, rhythm_text, Inches(0.5), Inches(4.9), Inches(5.8), Inches(2.2), 10, DARK_GRAY)

add_shape_box(slide9, Inches(6.5), Inches(4.3), Inches(6), Inches(0.5), GREEN, "复制与规模化", WHITE, 13)
scale_text = (
    "第1个试点：验证服务包模板和交付流程\n"
    "第2-3个试点：复制模式，缩短交付周期至30-35天\n"
    "第4个起：标准化服务包，FDE从2人降至1.5人\n"
    "12个月目标：完成10个试点，形成2个行业解决方案\n"
    "规模复制关键：连接器组件库、Agent模板库、培训体系"
)
add_text_box(slide9, scale_text, Inches(6.5), Inches(4.9), Inches(6), Inches(2.2), 10, DARK_GRAY)

# ============ Slide 10: Usage Boundaries ============
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10)
add_title_box(slide10, "使用边界与操作类型")
add_subtitle(slide10, "明确浏览器操作、电脑操作、API/连接器、人工审批与客户授权的边界")

boundary_data = [
    ["操作类型", "适用场景", "授权要求", "风险等级", "示例"],
    ["浏览器操作", "网页研究、竞品分析、公开数据采集", "无需客户授权（公开数据）", "低", "搜索行业报告、抓取公开价格"],
    ["电脑操作", "文档生成、数据分析、演示制作", "客户授权指定电脑/系统", "中", "生成采购报告、制作汇报PPT"],
    ["API/连接器", "回写业务系统、数据同步、流程触发", "客户IT部门授权+API密钥", "高", "写入ERP采购单、同步CRM数据"],
    ["人工审批", "涉及资金、合同、决策的最终确认", "业务负责人审批+签字", "极高", "采购审批、合同签署、付款"],
    ["客户授权", "数据访问、系统连接、Agent部署", "书面授权+安全协议", "极高", "系统访问、数据读取、Agent上线"],
]
add_table(slide10, boundary_data, Inches(0.3), Inches(1.6), Inches(12.7), Inches(4),
          col_widths=[Inches(1.5), Inches(3), Inches(2.5), Inches(1.2), Inches(4.5)])

add_text_box(slide10, 
    "原则：FieldPilot AI 的所有操作遵循「最小授权 + 可审计 + 可中止」原则。高风险操作必须经人工审批，不得自动执行。\n"
    "合规红线：不处理受监管数据（金融交易、医疗诊断）；不代客决策；AI输出必须经人工审核后方可进入业务系统。",
    Inches(0.3), Inches(5.8), Inches(12.5), Inches(1.2), 11, DARK_GRAY, True)

# ============ Slide 11: Compliance & Risk ============
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide11)
add_title_box(slide11, "合规与风险控制")
add_subtitle(slide11, "受监管行业说明 | 数据安全 | 责任边界")

add_shape_box(slide11, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.5), ACCENT_BLUE, "数据安全与合规", WHITE, 13)
sec_text = (
    "1. 数据隔离：客户数据驻留客户环境，FieldPilot不存储原始业务数据\n"
    "2. 传输加密：所有API通信使用TLS 1.2+，连接器使用客户密钥\n"
    "3. 审计日志：所有Agent操作可追溯，日志保留≥180天\n"
    "4. 权限治理：基于角色的访问控制（RBAC），最小权限原则\n"
    "5. 数据分类：按敏感度分级，高敏感数据不出域\n"
    "6. 合规对接：遵循《数据安全法》《个人信息保护法》要求"
)
add_text_box(slide11, sec_text, Inches(0.5), Inches(2.2), Inches(5.8), Inches(4), 10, DARK_GRAY)

add_shape_box(slide11, Inches(6.5), Inches(1.6), Inches(6), Inches(0.5), RED, "受监管行业边界", WHITE, 13)
reg_text = (
    "金融服务（暂缓）：\n"
    "• 高风险应用须风控委员会审批（金融监管总局指导意见）\n"
    "• 个人信息不得用于模型训练\n"
    "• 智能体须「辅助引导，不代客决策」\n"
    "• 需网信部门备案\n\n"
    "医疗健康（暂缓）：\n"
    "• AI医疗器械须注册审批（国家药监局）\n"
    "• 医疗数据属敏感个人信息\n"
    "• 「人机协同，以人为主」原则\n"
    "• 试点周期预计超12个月，不适合45天试点模式"
)
add_text_box(slide11, reg_text, Inches(6.5), Inches(2.2), Inches(6), Inches(4.5), 10, DARK_GRAY)

# ============ Slide 12: Summary & Next Steps ============
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide12, DARK_BLUE)

shape = slide12.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11), Inches(1))
tf = shape.text_frame
p = tf.paragraphs[0]
p.text = "总结与下一步"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = 'Microsoft YaHei'
p.alignment = PP_ALIGN.CENTER

next_steps = [
    ("1", "确认优先行业", "管理层审批制造业+零售电商为优先行业，暂缓金融与医疗"),
    ("2", "启动首个试点", "从10家目标企业中选定1家启动45天试点（建议制造业中型企业）"),
    ("3", "配置交付团队", "指派2名FDE+1名平台工程师+1名销售组成首个交付小组"),
    ("4", "建立服务包模板", "标准化FDE服务包、连接器组件库和Agent模板库"),
    ("5", "90天里程碑评审", "90天后评审3个试点进展、500万管线达成率和毛利率"),
]
for i, (num, title, desc) in enumerate(next_steps):
    top = Inches(2 + i * 0.95)
    shape = slide12.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), top, Inches(0.6), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x4C, 0xAF, 0x50)
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    add_text_box(slide12, title, Inches(1.8), top, Inches(4), Inches(0.4), 14, WHITE, True)
    add_text_box(slide12, desc, Inches(1.8), top + Inches(0.4), Inches(10), Inches(0.4), 11, RGBColor(0xB4, 0xC7, 0xE7))

# Footer
shape = slide12.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11), Inches(0.5))
tf = shape.text_frame
p = tf.paragraphs[0]
p.text = "FieldPilot AI — FDE 商业化方案  |  2026年8月  |  详见 01-04 附件"
p.font.size = Pt(10)
p.font.color.rgb = RGBColor(0x8D, 0xA0, 0xB5)
p.font.name = 'Microsoft YaHei'
p.alignment = PP_ALIGN.CENTER

# Save
output_path = Path(__file__).resolve().parent / '03-fde-commercialization-plan.pptx'
prs.save(output_path)
print(f'PPTX saved to: {output_path}')
size = output_path.stat().st_size
print(f'File size: {size} bytes')
print(f'Slides: {len(prs.slides)}')
